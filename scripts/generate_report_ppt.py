from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("学院学生综合服务与党团管理平台_结项汇报.pptx")

FONT = "Microsoft YaHei"
FONT_EN = "Aptos"

MAROON = RGBColor(125, 28, 55)
MAROON_DARK = RGBColor(88, 20, 39)
GOLD = RGBColor(184, 133, 54)
CREAM = RGBColor(248, 245, 238)
PAPER = RGBColor(255, 253, 248)
INK = RGBColor(42, 39, 36)
MUTED = RGBColor(105, 97, 89)
GREEN = RGBColor(46, 137, 86)
BLUE = RGBColor(48, 105, 151)
ORANGE = RGBColor(194, 104, 44)
GRAY_LINE = RGBColor(222, 216, 204)
WHITE = RGBColor(255, 255, 255)


def set_font(run, size=18, bold=False, color=INK, name=FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_shape_text(shape, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color)
    return shape


def add_box(slide, x, y, w, h, fill=PAPER, line=GRAY_LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1)
    return box


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return set_shape_text(shape, text, size=size, bold=bold, color=color, align=align)


def add_pill(slide, x, y, w, h, text, fill=MAROON, color=WHITE, size=12):
    pill = add_box(slide, x, y, w, h, fill=fill, line=fill, radius=True)
    set_shape_text(pill, text, size=size, bold=True, color=color, align=PP_ALIGN.CENTER)
    pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return pill


def add_footer(slide, idx):
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.55),
        Inches(7.08),
        Inches(12.78),
        Inches(7.08),
    ).line.color.rgb = RGBColor(226, 220, 209)
    add_text(slide, 0.62, 7.16, 7.0, 0.22, "学院学生综合服务与党团管理平台 · 结项汇报", 8, False, MUTED)
    add_text(slide, 12.1, 7.16, 0.65, 0.22, f"{idx:02d}", 8, True, MAROON, PP_ALIGN.RIGHT)


def add_title(slide, title, kicker=None, idx=None):
    if kicker:
        add_pill(slide, 0.65, 0.35, 1.35, 0.32, kicker, fill=MAROON)
    add_text(slide, 0.65, 0.78, 9.0, 0.52, title, 27, True, MAROON_DARK)
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.65),
        Inches(1.42),
        Inches(12.1),
        Inches(1.42),
    ).line.color.rgb = GRAY_LINE
    if idx is not None:
        add_footer(slide, idx)


def add_bullets(slide, x, y, w, h, items, size=15, color=INK, bullet_color=MAROON):
    top = y
    for item in items:
        add_pill(slide, x, top + 0.05, 0.18, 0.18, "", fill=bullet_color)
        add_text(slide, x + 0.32, top, w - 0.32, 0.42, item, size, False, color)
        top += 0.48


def add_metric(slide, x, y, w, h, value, label, fill=MAROON, sub=None):
    card = add_box(slide, x, y, w, h, fill=PAPER, line=RGBColor(225, 214, 199))
    add_text(slide, x + 0.18, y + 0.16, w - 0.36, 0.5, value, 28, True, fill, PP_ALIGN.CENTER)
    add_text(slide, x + 0.18, y + 0.77, w - 0.36, 0.32, label, 12, True, INK, PP_ALIGN.CENTER)
    if sub:
        add_text(slide, x + 0.18, y + 1.11, w - 0.36, 0.28, sub, 9, False, MUTED, PP_ALIGN.CENTER)
    return card


def add_status_chip(slide, x, y, text, status):
    colors = {
        "full": (GREEN, "完全"),
        "core": (BLUE, "核心"),
        "partial": (ORANGE, "部分"),
        "base": (GOLD, "基本"),
    }
    fill, label = colors[status]
    add_pill(slide, x, y, 0.72, 0.24, label, fill=fill, size=9)
    add_text(slide, x + 0.82, y - 0.01, 2.75, 0.28, text, 10.5, False, INK)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM

    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    hero.fill.solid()
    hero.fill.fore_color.rgb = CREAM
    hero.line.fill.background()

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.06))
    band.fill.solid()
    band.fill.fore_color.rgb = MAROON
    band.line.fill.background()

    add_text(slide, 0.65, 0.33, 5.0, 0.28, "《软件工程导论》课程结项汇报", 13, True, WHITE)
    add_pill(slide, 10.9, 0.31, 1.75, 0.32, "8 分钟展示版", fill=GOLD, size=11)

    add_text(slide, 0.78, 1.72, 9.7, 0.56, "学院学生综合服务与党团管理平台", 34, True, MAROON_DARK)
    add_text(slide, 0.82, 2.44, 8.4, 0.38, "一站式学生服务 · 党团流程管理 · 智能政策问答 · 精准通知与审批", 16, False, INK)

    cards = [
        ("Web 管理端", "Vue 3 + Vite"),
        ("后端 MVP", "Spring Boot"),
        ("智能问答", "RAG / Codex-first"),
        ("自动化测试", "104 cases passed"),
    ]
    x = 0.82
    for title, desc in cards:
        add_box(slide, x, 3.48, 2.55, 1.08, fill=PAPER, line=RGBColor(224, 214, 197))
        add_text(slide, x + 0.2, 3.68, 2.16, 0.28, title, 15, True, MAROON_DARK, PP_ALIGN.CENTER)
        add_text(slide, x + 0.2, 4.05, 2.16, 0.25, desc, 10.5, False, MUTED, PP_ALIGN.CENTER)
        x += 2.8

    add_text(slide, 0.82, 6.2, 8.5, 0.28, "团队成员：孙谦，田家铭，刘凡硕，林靖涛", 12, False, INK)
    add_text(slide, 0.82, 6.58, 8.5, 0.28, "日期：2026 年 6 月 30 日", 12, False, MUTED)


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "汇报路线", "AGENDA", 2)
    items = [
        ("01", "项目概述", "解决学院学生服务和党团管理的分散问题"),
        ("02", "功能完成", "核心闭环与扩展能力当前状态"),
        ("03", "设计方案", "前后端分离、模块化单体、权限和数据范围"),
        ("04", "进度协作", "从需求原型到全栈 MVP 的迭代路径"),
        ("05", "开发测试", "关键功能、自动化测试、缺陷修复"),
        ("06", "工具与总结", "大模型使用、问题经验、后续计划"),
    ]
    y = 1.78
    for no, title, desc in items:
        add_pill(slide, 0.9, y + 0.02, 0.58, 0.34, no, fill=MAROON, size=12)
        add_text(slide, 1.7, y - 0.02, 2.1, 0.32, title, 16, True, MAROON_DARK)
        add_text(slide, 3.78, y, 7.3, 0.3, desc, 13, False, INK)
        y += 0.72


def slide_problem_scope(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "项目定位：把分散事务收敛到统一平台", "OVERVIEW", 3)
    add_text(slide, 0.72, 1.68, 5.5, 0.4, "服务对象", 18, True, MAROON_DARK)
    roles = [
        ("普通学生", "个人画像、通知、党团进度、申请办理"),
        ("班团骨干", "授权范围内成员进展与催办"),
        ("管理老师", "学生画像、知识库、通知、审批、日志"),
        ("学院领导", "学院汇总看板与运行态势"),
    ]
    y = 2.22
    for role, desc in roles:
        add_box(slide, 0.78, y, 5.4, 0.64, fill=PAPER, line=RGBColor(226, 216, 202))
        add_text(slide, 1.02, y + 0.14, 1.35, 0.26, role, 13, True, MAROON_DARK)
        add_text(slide, 2.5, y + 0.14, 3.28, 0.26, desc, 10.5, False, INK)
        y += 0.78

    add_text(slide, 7.1, 1.68, 4.8, 0.4, "核心目标", 18, True, MAROON_DARK)
    goals = [
        "学生事务线上化",
        "政策咨询智能化",
        "党团流程可视化",
        "通知推送精准化",
        "业务操作可追溯",
    ]
    add_bullets(slide, 7.15, 2.24, 4.85, 2.7, goals, size=15)
    add_metric(slide, 7.12, 5.05, 1.45, 1.1, "4", "核心角色", MAROON)
    add_metric(slide, 8.85, 5.05, 1.45, 1.1, "10+", "业务模块", BLUE)
    add_metric(slide, 10.58, 5.05, 1.45, 1.1, "1200", "目标学生规模", GOLD)


def slide_deliverables(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "当前交付物：从原型推进到可演示全栈项目", "DELIVERY", 4)
    deliverables = [
        ("需求与设计文档", "README / architecture / API / DB / 权限 / 字典 / 设计文档"),
        ("正式 Web 前端", "Vue 3 + Vite + Router + Pinia + Axios，学生端与管理端"),
        ("后端 MVP", "Spring Boot REST API，认证、学生、知识库、党团、通知、申请等模块"),
        ("本地数据库", "H2 开发与测试库，结构对齐 Kingbase/PostgreSQL 兼容设计"),
        ("移动端演示", "微信小程序学生端基础版，覆盖高频入口"),
        ("测试与部署", "Surefire 测试报告、Vite 构建、运行包和校园网部署脚本"),
    ]
    x_positions = [0.76, 4.72, 8.68]
    y_positions = [1.8, 4.0]
    for idx, (title, desc) in enumerate(deliverables):
        x = x_positions[idx % 3]
        y = y_positions[idx // 3]
        add_box(slide, x, y, 3.15, 1.58, fill=PAPER, line=RGBColor(226, 216, 202))
        add_text(slide, x + 0.22, y + 0.18, 2.65, 0.34, title, 15, True, MAROON_DARK)
        add_text(slide, x + 0.22, y + 0.64, 2.65, 0.68, desc, 10.5, False, INK)


def slide_feature_status(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "功能满足性：核心业务闭环已具备演示能力", "FEATURES", 5)
    left = [
        ("多角色登录与权限", "full"),
        ("学生端首页", "full"),
        ("管理端首页", "full"),
        ("院内申请与审批", "core"),
        ("知识库与智能问答", "core"),
        ("文件上传下载", "base"),
    ]
    right = [
        ("学生画像与信息管理", "partial"),
        ("党团事务流程", "base"),
        ("精准通知", "partial"),
        ("奖励荣誉展示", "partial"),
        ("审计日志与系统日志", "partial"),
        ("微信小程序演示版", "partial"),
    ]
    add_box(slide, 0.78, 1.78, 5.72, 4.55, fill=PAPER, line=RGBColor(226, 216, 202))
    add_box(slide, 6.83, 1.78, 5.72, 4.55, fill=PAPER, line=RGBColor(226, 216, 202))
    add_text(slide, 1.05, 2.02, 4.8, 0.32, "已完成 / 核心闭环", 16, True, MAROON_DARK)
    add_text(slide, 7.1, 2.02, 4.8, 0.32, "部分实现 / 后续增强", 16, True, MAROON_DARK)
    y = 2.62
    for text, status in left:
        add_status_chip(slide, 1.08, y, text, status)
        y += 0.52
    y = 2.62
    for text, status in right:
        add_status_chip(slide, 7.13, y, text, status)
        y += 0.52
    add_text(slide, 0.92, 6.45, 11.5, 0.35, "说明：完整后台版本管理、真实邮件/微信发送、Kingbase 生产迁移和日志自动落库属于后续增强项。", 11, False, MUTED)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "总体架构：前后端分离 + 模块化单体", "ARCHITECTURE", 6)
    layers = [
        (0.82, 1.8, 11.65, 0.82, "用户层", "普通学生 / 班团骨干 / 管理老师 / 学院领导 / 系统管理员", MAROON),
        (0.82, 2.95, 11.65, 0.82, "前端层", "Vue Web 前端 · 路由守卫 · Pinia 登录态 · Axios API Client · 小程序演示端", BLUE),
        (0.82, 4.1, 11.65, 0.82, "后端层", "Spring Boot REST API (/api/v1) · 统一响应 · 鉴权过滤器 · 全局异常处理", GREEN),
        (0.82, 5.25, 11.65, 0.82, "业务模块", "认证 / 学生画像 / 知识库 / 党团 / 通知 / 申请 / 文件 / 荣誉 / 日志", GOLD),
        (0.82, 6.4, 11.65, 0.52, "数据与部署", "H2 开发测试库 · Kingbase 兼容设计 · 文件资源 · Nginx + systemd 部署脚本", MUTED),
    ]
    for x, y, w, h, tag, text, color in layers:
        add_box(slide, x, y, w, h, fill=PAPER, line=RGBColor(226, 216, 202))
        add_pill(slide, x + 0.2, y + 0.23, 1.28, 0.32, tag, fill=color, size=10)
        add_text(slide, x + 1.7, y + 0.25, w - 2.0, 0.28, text, 13, False, INK)


def slide_rag(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "核心功能示例一：知识库与 RAG 智能问答", "CORE 1", 7)
    steps = [
        ("学生提问", "支持连续追问和历史上下文"),
        ("知识检索", "优先检索已发布知识条目和模板附件"),
        ("RAG 回答", "Codex-first，失败降级 Responses API"),
        ("兜底策略", "关键词匹配 / 低置信度提示 / 未检索到依据"),
        ("来源展示", "返回 sources，学生可追溯答案依据"),
    ]
    x = 0.72
    for idx, (title, desc) in enumerate(steps):
        add_box(slide, x, 2.08, 2.18, 1.55, fill=PAPER, line=RGBColor(226, 216, 202))
        add_pill(slide, x + 0.18, 2.28, 0.5, 0.28, str(idx + 1), fill=MAROON, size=10)
        add_text(slide, x + 0.8, 2.22, 1.18, 0.3, title, 13, True, MAROON_DARK)
        add_text(slide, x + 0.2, 2.75, 1.74, 0.48, desc, 9.5, False, INK, PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 2.18), Inches(2.86), Inches(x + 2.5), Inches(2.86))
            conn.line.color.rgb = GOLD
            conn.line.width = Pt(2)
        x += 2.5
    add_box(slide, 0.86, 4.46, 11.45, 1.25, fill=RGBColor(255, 250, 239), line=RGBColor(230, 217, 188))
    add_text(slide, 1.13, 4.68, 10.7, 0.3, "设计原则：有来源优先，无依据不编造", 18, True, MAROON_DARK, PP_ALIGN.CENTER)
    add_text(slide, 1.25, 5.15, 10.42, 0.28, "当 Codex、向量库或模型接口不可用时，系统自动降级，保证演示主流程稳定。", 12.5, False, INK, PP_ALIGN.CENTER)


def slide_application(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "核心功能示例二：院内申请与审批状态机", "CORE 2", 8)
    states = [
        ("提交申请", "submitted"),
        ("审核中", "reviewing"),
        ("审批通过", "approved"),
        ("审批驳回", "rejected"),
        ("申请撤回", "revoked"),
    ]
    coords = [(0.92, 2.15), (3.1, 2.15), (5.42, 1.45), (5.42, 2.86), (3.1, 4.35)]
    for (title, code), (x, y) in zip(states, coords):
        color = GREEN if code == "approved" else ORANGE if code in {"rejected", "revoked"} else MAROON
        add_box(slide, x, y, 1.75, 0.84, fill=PAPER, line=color)
        add_text(slide, x + 0.1, y + 0.16, 1.55, 0.22, title, 12.5, True, color, PP_ALIGN.CENTER)
        add_text(slide, x + 0.1, y + 0.48, 1.55, 0.18, code, 8.5, False, MUTED, PP_ALIGN.CENTER)

    connectors = [
        ((2.67, 2.57), (3.1, 2.57)),
        ((4.85, 2.57), (5.42, 1.86)),
        ((4.85, 2.57), (5.42, 3.27)),
        ((3.95, 2.99), (3.95, 4.35)),
    ]
    for (x1, y1), (x2, y2) in connectors:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = GOLD
        conn.line.width = Pt(2)

    add_box(slide, 8.0, 1.72, 4.25, 3.95, fill=PAPER, line=RGBColor(226, 216, 202))
    add_text(slide, 8.28, 1.98, 3.7, 0.3, "后端约束", 17, True, MAROON_DARK)
    bullets = [
        "状态只能由审批接口驱动",
        "非法流转返回 40900",
        "状态更新 + 审批记录同事务",
        "管理老师按班级数据范围审批",
        "越权、重复审批均有测试覆盖",
    ]
    add_bullets(slide, 8.32, 2.55, 3.55, 2.55, bullets, size=11.5, bullet_color=GREEN)


def slide_ux(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "用户体验：学生端轻量办理，管理端高效处理", "UX", 9)
    add_box(slide, 0.78, 1.82, 5.55, 4.75, fill=PAPER, line=RGBColor(226, 216, 202))
    add_box(slide, 6.95, 1.82, 5.55, 4.75, fill=PAPER, line=RGBColor(226, 216, 202))
    add_text(slide, 1.1, 2.12, 4.7, 0.34, "学生端", 18, True, MAROON_DARK, PP_ALIGN.CENTER)
    add_text(slide, 7.28, 2.12, 4.7, 0.34, "管理端", 18, True, MAROON_DARK, PP_ALIGN.CENTER)
    add_bullets(slide, 1.15, 2.88, 4.6, 2.2, [
        "首页聚合待办、通知、党团阶段",
        "知识库提问展示来源与置信度",
        "申请、通知、材料提交集中办理",
        "优先适配移动端与窄屏场景",
    ], size=13)
    add_bullets(slide, 7.32, 2.88, 4.6, 2.2, [
        "侧边栏组织多业务模块",
        "列表支持筛选、分页、状态标签",
        "审批与敏感信息强调留痕",
        "优先满足 PC 批量管理效率",
    ], size=13, bullet_color=BLUE)
    add_text(slide, 1.1, 5.8, 4.7, 0.25, "少操作、快办理", 14, True, GREEN, PP_ALIGN.CENTER)
    add_text(slide, 7.28, 5.8, 4.7, 0.25, "高密度、可追踪", 14, True, BLUE, PP_ALIGN.CENTER)


def slide_progress(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "进度安排：文档先行，逐步推进到全栈 MVP", "TIMELINE", 10)
    milestones = [
        ("04/19", "需求整理\n静态 Web 原型"),
        ("05/07", "小程序\n学生端演示"),
        ("05/11", "软件设计\n文档完成"),
        ("05/18", "Vue 前端\nSpring 后端"),
        ("05/27", "注册改密\n画像维护"),
        ("05/28", "申请提交\n数据范围修复"),
        ("06", "RAG 问答\n测试与部署"),
    ]
    base_y = 3.55
    start_x = 0.95
    step = 1.78
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.12), Inches(base_y), Inches(11.75), Inches(base_y)).line.color.rgb = GOLD
    for i, (date, text) in enumerate(milestones):
        x = start_x + i * step
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(base_y - 0.16), Inches(0.32), Inches(0.32))
        dot.fill.solid()
        dot.fill.fore_color.rgb = MAROON if i in {0, 3, 6} else GOLD
        dot.line.fill.background()
        add_text(slide, x - 0.25, 2.65 if i % 2 == 0 else 4.02, 0.85, 0.24, date, 12, True, MAROON_DARK, PP_ALIGN.CENTER)
        add_text(slide, x - 0.55, 2.95 if i % 2 == 0 else 4.32, 1.45, 0.58, text, 9.5, False, INK, PP_ALIGN.CENTER)
    add_box(slide, 1.02, 5.78, 10.9, 0.58, fill=RGBColor(255, 250, 239), line=RGBColor(230, 217, 188))
    add_text(slide, 1.25, 5.92, 10.4, 0.22, "实际结果：从原型与文档交付，扩展为正式前端工程、后端 MVP、自动化测试和部署脚本。", 12.5, True, MAROON_DARK, PP_ALIGN.CENTER)


def slide_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "团队协作：按模块推进，用文档和测试收口", "TEAM", 11)
    members = [
        ("孙谦", "整体推进 / 后端接口 / RAG / 测试验证"),
        ("田家铭", "部署运行 / 脚本 / 环境验证 / 联调支持"),
        ("刘凡硕", "软件设计文档 / 小程序 / 需求整理"),
        ("林靖涛", "前端完善 / 测试用例 / 缺陷验证 / 汇报支持"),
    ]
    x = 0.82
    for name, role in members:
        add_box(slide, x, 1.86, 2.82, 1.36, fill=PAPER, line=RGBColor(226, 216, 202))
        add_text(slide, x + 0.2, 2.08, 2.42, 0.28, name, 18, True, MAROON_DARK, PP_ALIGN.CENTER)
        add_text(slide, x + 0.24, 2.5, 2.34, 0.42, role, 9.8, False, INK, PP_ALIGN.CENTER)
        x += 3.0
    add_box(slide, 0.95, 4.04, 11.2, 1.65, fill=PAPER, line=RGBColor(226, 216, 202))
    add_text(slide, 1.24, 4.28, 10.55, 0.28, "协作工具与流程", 17, True, MAROON_DARK)
    add_bullets(slide, 1.28, 4.82, 10.1, 0.8, [
        "Git/GitHub 管理代码与版本记录；docs/ 沉淀架构、接口、数据库、权限与字典约束",
        "developer-diary 与 problems-track 记录阶段迭代、问题反馈和修复点",
    ], size=12.5, bullet_color=BLUE)


def slide_dev_highlights(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "开发过程亮点：工程化、权限化、可部署", "DEVELOPMENT", 12)
    highlights = [
        ("工程化前端", "路由守卫、Pinia 登录态、API client、通用组件、mock 层"),
        ("后端基础设施", "统一响应、requestId、异常处理、鉴权过滤器、H2 初始化"),
        ("权限与数据范围", "学生本人、班级范围、管理角色、学院领导范围分层控制"),
        ("部署交付", "运行包、Nginx、systemd、校园网一键脚本、状态检查"),
    ]
    y = 1.8
    for title, desc in highlights:
        add_box(slide, 0.86, y, 11.45, 0.86, fill=PAPER, line=RGBColor(226, 216, 202))
        add_pill(slide, 1.12, y + 0.25, 1.42, 0.3, title, fill=MAROON if y < 3 else BLUE, size=10)
        add_text(slide, 2.78, y + 0.26, 8.8, 0.26, desc, 13, False, INK)
        y += 1.05


def slide_testing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "测试结果：自动化验证通过，核心风险有覆盖", "TESTING", 13)
    add_metric(slide, 0.88, 1.86, 2.25, 1.5, "104", "后端测试用例", MAROON, "JUnit / Spring Boot Test")
    add_metric(slide, 3.55, 1.86, 2.25, 1.5, "14", "Surefire 报告", BLUE, "模块测试与集成测试")
    add_metric(slide, 6.22, 1.86, 2.25, 1.5, "0", "failures / errors", GREEN, "无失败、无错误")
    add_metric(slide, 8.89, 1.86, 2.25, 1.5, "Pass", "前端构建与测试", GOLD, "Vite build / kbQa")
    add_box(slide, 0.88, 4.0, 5.55, 1.85, fill=PAPER, line=RGBColor(226, 216, 202))
    add_text(slide, 1.18, 4.26, 4.95, 0.3, "覆盖重点", 16, True, MAROON_DARK)
    add_bullets(slide, 1.22, 4.75, 4.8, 0.85, [
        "认证、仪表盘、学生画像、党团、通知、申请、文件",
        "知识库、RAG 客户端、荣誉展示、数据范围和异常分支",
    ], size=11.5, bullet_color=GREEN)
    add_box(slide, 6.85, 4.0, 5.05, 1.85, fill=PAPER, line=RGBColor(226, 216, 202))
    add_text(slide, 7.15, 4.26, 4.45, 0.3, "典型修复", 16, True, MAROON_DARK)
    add_bullets(slide, 7.18, 4.75, 4.25, 0.85, [
        "审批列表 SQL 字段歧义与班级过滤",
        "知识库详情、申请提交、老师上传权限、RAG 降级",
    ], size=11.5, bullet_color=ORANGE)


def slide_ai_tools(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "工具使用：大模型加速，但关键规则人工把关", "TOOLS", 14)
    add_box(slide, 0.9, 1.78, 5.45, 4.45, fill=PAPER, line=RGBColor(226, 216, 202))
    add_box(slide, 6.92, 1.78, 5.45, 4.45, fill=PAPER, line=RGBColor(226, 216, 202))
    add_text(slide, 1.22, 2.06, 4.75, 0.3, "大模型主要用途", 17, True, MAROON_DARK)
    add_bullets(slide, 1.26, 2.58, 4.65, 2.5, [
        "需求提取、用户故事和接口草案",
        "Vue 页面、DTO、Service、测试骨架生成",
        "SQL / 构建 / 权限问题调试分析",
        "README、接口说明、部署说明、结项材料整理",
    ], size=12.5)
    add_text(slide, 7.24, 2.06, 4.75, 0.3, "人工把关边界", 17, True, MAROON_DARK)
    add_bullets(slide, 7.28, 2.58, 4.65, 2.5, [
        "权限、状态机、数据范围必须人工确认",
        "生成字段和接口必须对照仓库实际代码",
        "代码需要运行测试和构建验证",
        "团队分工、验收口径需按真实情况调整",
    ], size=12.5, bullet_color=BLUE)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    add_title(slide, "总结与后续计划", "SUMMARY", 15)
    add_box(slide, 0.86, 1.78, 5.55, 4.82, fill=PAPER, line=RGBColor(226, 216, 202))
    add_box(slide, 6.82, 1.78, 5.55, 4.82, fill=PAPER, line=RGBColor(226, 216, 202))
    add_text(slide, 1.18, 2.06, 4.95, 0.32, "本次结项成果", 18, True, MAROON_DARK)
    add_bullets(slide, 1.22, 2.62, 4.75, 2.5, [
        "完成从需求、设计、原型到全栈 MVP 的闭环",
        "学生端、管理端、后端接口和测试体系均可展示",
        "核心流程：登录、知识库问答、党团材料、通知、申请审批",
        "用自动化测试和部署脚本提升验收稳定性",
    ], size=12.5, bullet_color=GREEN)
    add_text(slide, 7.14, 2.06, 4.95, 0.32, "后续增强方向", 18, True, MAROON_DARK)
    add_bullets(slide, 7.18, 2.62, 4.75, 2.5, [
        "迁移 Kingbase 生产数据库并补迁移脚本",
        "完善知识库版本管理、通知真实渠道和日志自动落库",
        "补浏览器 E2E、小程序接口联调和依赖安全治理",
        "完善生产部署：HTTPS、备份、日志轮转、文件存储",
    ], size=12.5, bullet_color=BLUE)
    add_pill(slide, 3.9, 6.78, 5.55, 0.38, "谢谢，请老师和同学批评指正", fill=MAROON, size=14)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_agenda(prs)
    slide_problem_scope(prs)
    slide_deliverables(prs)
    slide_feature_status(prs)
    slide_architecture(prs)
    slide_rag(prs)
    slide_application(prs)
    slide_ux(prs)
    slide_progress(prs)
    slide_team(prs)
    slide_dev_highlights(prs)
    slide_testing(prs)
    slide_ai_tools(prs)
    slide_summary(prs)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
